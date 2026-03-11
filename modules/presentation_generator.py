"""
Generatore di presentazioni
"""

import json
from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .ollama_client import OllamaClient
from .document_parser import DocumentParser
from .image_extractor import ImageExtractor
import config


class PresentationGenerator:
    """Genera presentazioni da documenti"""
    
    def __init__(self, ollama_client: OllamaClient):
        self.ollama_client = ollama_client
        self.parser = DocumentParser()
        self.image_extractor = ImageExtractor(config.IMAGE_DIR)
    
    def generate_presentation_content(
        self,
        document_path: Path,
        model: str = config.DEFAULT_MODEL
    ) -> Dict:
        """
        Genera il contenuto della presentazione usando Ollama
        
        Returns:
            Dict con struttura presentazione
        """
        # 1. Parse documento
        parsed_doc = self.parser.parse(document_path)
        
        # 2. Estrai immagini
        images = self.image_extractor.extract_images(document_path)
        image_list = [img['filename'] for img in images]
        
        # 3. Calcola parametri
        word_count = self.parser.count_words(parsed_doc['content'])
        
        # 4. Costruisci prompt
        prompt = config.PRESENTATION_PROMPT_TEMPLATE.format(
            duration=config.TARGET_PRESENTATION_MINUTES,
            word_count=config.TARGET_WORD_COUNT,
            num_slides=config.ESTIMATED_SLIDES,
            content=parsed_doc['content'][:5000],  # Limita per context window
            images=", ".join(image_list) if image_list else "Nessuna immagine disponibile"
        )
        
        # 5. Genera con Ollama
        messages = [
            {"role": "system", "content": config.SYSTEM_PROMPT},
            {"role": "user", "content": prompt}
        ]
        
        response = self.ollama_client.chat(
            messages=messages,
            model=model,
            format='json'
        )
        
        # 6. Parse risposta JSON
        try:
            if hasattr(response, 'message'):
                content = response.message.content
            else:
                content = response.get('message', {}).get('content', '{}')
            
            presentation_data = json.loads(content)
            
            # Aggiungi metadata
            presentation_data['metadata'] = {
                'source_document': document_path.name,
                'source_format': parsed_doc['metadata']['format'],
                'images': images,
                'word_count': word_count
            }
            
            return presentation_data
        
        except json.JSONDecodeError as e:
            print(f"Errore parsing JSON: {e}")
            return {
                "error": "Impossibile parsare la risposta del modello",
                "raw_response": content if 'content' in locals() else str(response)
            }
    
    def create_pptx(
        self,
        presentation_data: Dict,
        output_path: Path
    ) -> Path:
        """
        Crea file PPTX da dati presentazione
        
        Returns:
            Path del file creato
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide titolo
        self._add_title_slide(
            prs,
            presentation_data.get('title', 'Presentazione'),
            presentation_data.get('subtitle', ''),
            presentation_data.get('author', '')
        )
        
        # Slides contenuto
        for slide_data in presentation_data.get('slides', []):
            slide_type = slide_data.get('type', 'content')
            
            if slide_type == 'title':
                self._add_title_slide(
                    prs,
                    slide_data.get('title', ''),
                    slide_data.get('content', [''])[0] if slide_data.get('content') else ''
                )
            elif slide_type == 'image' and slide_data.get('image'):
                self._add_image_slide(
                    prs,
                    slide_data.get('title', ''),
                    slide_data.get('content', []),
                    slide_data.get('image')
                )
            else:
                self._add_content_slide(
                    prs,
                    slide_data.get('title', ''),
                    slide_data.get('content', [])
                )
        
        # Salva
        prs.save(str(output_path))
        return output_path
    
    def _add_title_slide(
        self,
        prs: Presentation,
        title: str,
        subtitle: str = '',
        author: str = ''
    ):
        """Aggiunge slide titolo"""
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        
        if subtitle or author:
            subtitle_text = f"{subtitle}\n{author}" if subtitle and author else (subtitle or author)
            subtitle_placeholder.text = subtitle_text
    
    def _add_content_slide(
        self,
        prs: Presentation,
        title: str,
        content: List[str]
    ):
        """Aggiunge slide con bullet points"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for i, point in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
    
    def _add_image_slide(
        self,
        prs: Presentation,
        title: str,
        content: List[str],
        image_filename: str
    ):
        """Aggiunge slide con immagine"""
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Aggiungi titolo manualmente
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Aggiungi immagine
        image_path = config.IMAGE_DIR / image_filename
        if image_path.exists():
            left = Inches(1)
            top = Inches(2)
            slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=Inches(8)
            )
